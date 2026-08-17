import os
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
REPORT_DIR = os.path.join(BASE, "report")
IMG_DIR = os.path.join(REPORT_DIR, "images")
PPT_PATH = os.path.join(REPORT_DIR, "presentation.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]

# Theme Colors
COLOR_BG = RGBColor(248, 250, 252)       # Slate 50
COLOR_PRIMARY = RGBColor(15, 23, 42)    # Slate 900
COLOR_SECONDARY = RGBColor(37, 99, 235) # Blue 600
COLOR_CARD_BG = RGBColor(255, 255, 255) # White
COLOR_TEXT_DARK = RGBColor(30, 41, 59)   # Slate 800
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
COLOR_ACCENT = RGBColor(16, 185, 129)   # Emerald 500

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="BFSI / AI-ML CAPSTONE PROJECT"):
    # Header Card Banner
    header_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.0)
    )
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = COLOR_PRIMARY
    header_box.line.fill.background()

    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.12)

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_SECONDARY

    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)

def add_card(slide, left, top, width, height, title, content_bullets):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    p_head = tf.paragraphs[0]
    p_head.text = title
    p_head.font.size = Pt(15)
    p_head.font.bold = True
    p_head.font.color.rgb = COLOR_PRIMARY
    p_head.space_after = Pt(10)

    for bullet in content_bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(6)
        p.level = 0

# --- SLIDE 1: Title Slide ---
slide1 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide1)

# Large Title Box
tbox = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
tbox.fill.solid()
tbox.fill.fore_color.rgb = COLOR_PRIMARY
tbox.line.fill.background()

tf1 = tbox.text_frame
tf1.margin_left = Inches(0.6)
tf1.margin_top = Inches(0.6)

p1 = tf1.paragraphs[0]
p1.text = "SYNTHETIC IDENTITY FRAUD DETECTION"
p1.font.size = Pt(32)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)

p2 = tf1.add_paragraph()
p2.text = "Digital Lending Onboarding via KYC Verification & Behavioral Signals"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = COLOR_SECONDARY
p2.space_after = Pt(30)

p3 = tf1.add_paragraph()
p3.text = "Final Year Capstone Project — BFSI / AI-ML Domain\nMachine Learning Pipeline & Live Interactive Web Dashboard"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(203, 213, 225)

# --- SLIDE 2: Executive Summary & Motivation ---
slide2 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide2)
add_header(slide2, "Executive Summary & Industry Problem")

add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
         "The Threat: Synthetic Identity Fraud", [
    "• Synthetic Identity Fraud (SIF) combines real PII (valid PAN) with fake details (burner phone, fake address).",
    "• Fastest-growing financial crime category in digital lending onboarding.",
    "• Bypasses traditional static KYC checks because individual field queries return clean responses.",
    "• No individual victim exists initially to raise an alert (credit bust-out risk)."
])

add_card(slide2, Inches(6.8), Inches(1.8), Inches(5.733), Inches(5.0),
         "Our Machine Learning Solution", [
    "• Multimodal Risk Fusion: Blends KYC consistency, identity freshness, behavioral biometrics, and device velocity.",
    "• Machine Learning Engine: Random Forest Classifier tuned with balanced class weights for imbalanced data.",
    "• Real-Time Scoring: Live risk tiering (<30% Low, 30-60% Medium, >60% High).",
    "• Streamlit Web App: Interactive onboarding decision interface."
])

# --- SLIDE 3: Feature Dictionary Architecture ---
slide3 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide3)
add_header(slide3, "Feature Engineering & 12 Core Signals")

add_card(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4),
         "1. KYC Document Consistency", [
    "• name_address_mismatch_score (Stated vs. record distance)",
    "• dob_pan_mismatch (Binary DOB mismatch against PAN)",
    "• document_reuse_count (Historical reuse of document image)"
])

add_card(slide3, Inches(6.8), Inches(1.8), Inches(5.733), Inches(2.4),
         "2. Identity Freshness Telemetry", [
    "• phone_age_days (Age of mobile line in days)",
    "• email_age_days (Domain/account age of email)",
    "• credit_bureau_hit (Binary existence of bureau file)",
    "• social_footprint_score (Public digital footprint index)"
])

add_card(slide3, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.3),
         "3. Behavioral Biometrics", [
    "• session_fill_time_sec (Total application duration)",
    "• typing_speed_variance (Keypress cadence variance)"
])

add_card(slide3, Inches(6.8), Inches(4.5), Inches(5.733), Inches(2.3),
         "4. Device & Network Velocity", [
    "• device_reuse_across_apps (Identities on same device)",
    "• application_velocity_24h (Applications from IP in 24h)",
    "• ip_geolocation_mismatch (IP location vs. address)"
])

# --- SLIDE 4: Machine Learning Methodology ---
slide4 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide4)
add_header(slide4, "Machine Learning Model Methodology")

add_card(slide4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0),
         "Balanced Random Forest Architecture", [
    "• Dataset Size: 8,000 synthetic onboarding records (80/20 Stratified Train/Test split).",
    "• Class Imbalance Handling: Applied class_weight='balanced' to handle 15% fraud class ratio.",
    "• Hyperparameter Specs: 300 Decision Trees, Max Depth = 8, Random Seed = 42.",
    "• Why Random Forest? Robust against non-linear interactions, non-parametric feature scaling, immune to extreme outliers.",
    "• Evaluation Strategy: Evaluated strictly on held-out test data using Precision, Recall, F1, and ROC-AUC."
])

# --- SLIDE 5: Model Results & Visual Evaluation ---
slide5 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide5)
add_header(slide5, "Model Performance Results & Key Metrics")

# Metric KPI Cards
kpi_card = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.0))
kpi_card.fill.solid()
kpi_card.fill.fore_color.rgb = COLOR_PRIMARY
kpi_card.line.fill.background()

tf_kpi = kpi_card.text_frame
p_k = tf_kpi.paragraphs[0]
p_k.text = "ROC-AUC: 0.8834   |   Fraud Precision: 0.9500   |   Fraud Recall: 0.7500   |   Fraud F1: 0.8400   |   Accuracy: 0.9600"
p_k.font.size = Pt(16)
p_k.font.bold = True
p_k.font.color.rgb = COLOR_ACCENT
p_k.alignment = PP_ALIGN.CENTER

img_roc = os.path.join(IMG_DIR, "roc_curve.png")
img_cm = os.path.join(IMG_DIR, "confusion_matrix.png")

if os.path.exists(img_roc):
    slide5.shapes.add_picture(img_roc, Inches(0.8), Inches(3.0), Inches(5.6), Inches(4.0))

if os.path.exists(img_cm):
    slide5.shapes.add_picture(img_cm, Inches(6.8), Inches(3.0), Inches(5.6), Inches(4.0))

# --- SLIDE 6: Feature Importances ---
slide6 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide6)
add_header(slide6, "Top Predictive Signals (Feature Importances)")

add_card(slide6, Inches(0.8), Inches(1.8), Inches(5.0), Inches(5.0),
         "Key Feature Findings", [
    "1. Email Age (0.232): Primary signal; freshly created emails (<10 days) strongly indicate synthetic identities.",
    "2. Phone Number Age (0.220): Burner phone numbers purchased recently correlate heavily with fraud.",
    "3. Name/Address Mismatch (0.182): Address distance score catches spliced identity fragments.",
    "4. Social Footprint (0.097): Lack of digital footprint reflects zero historical presence.",
    "5. Session Fill Time (0.096): Rushed or scripted completion indicates bot automation."
])

img_fi = os.path.join(IMG_DIR, "feature_importance.png")
if os.path.exists(img_fi):
    slide6.shapes.add_picture(img_fi, Inches(6.0), Inches(1.8), Inches(6.5), Inches(5.0))

# --- SLIDE 7: Solution Architecture & Risk Tiering ---
slide7 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide7)
add_header(slide7, "Solution Architecture & Decision Pipeline")

add_card(slide7, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
         "Low Risk Band (< 30%)", [
    "• Risk Score: < 0.30",
    "• Profile: Established phone/email age, clean credit bureau history, normal typing cadence.",
    "• Action: Automated Instant Approval & Loan Disbursal."
])

add_card(slide7, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
         "Medium Risk Band (30-60%)", [
    "• Risk Score: 0.30 – 0.60",
    "• Profile: Moderate address mismatch or minor device velocity anomaly.",
    "• Action: Step-Up Verification (Video KYC, OTP challenge)."
])

add_card(slide7, Inches(8.8), Inches(1.8), Inches(3.733), Inches(5.0),
         "High Risk Band (> 60%)", [
    "• Risk Score: > 0.60",
    "• Profile: High document reuse count, burner phone line, high application velocity.",
    "• Action: Manual Underwriter Review / Instant Block."
])

# --- SLIDE 8: Interactive Streamlit Dashboard Demo ---
slide8 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide8)
add_header(slide8, "Interactive Streamlit Web Dashboard")

add_card(slide8, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0),
         "Real-Time Onboarding Decision Portal (app/dashboard.py)", [
    "• Preset Profile Loader: Select 'Typical Legit Applicant' vs. 'Suspicious Applicant' with a single click.",
    "• Dynamic Signal Inputs: Sliders and numerical inputs for all 12 KYC & behavioral features.",
    "• Real-Time Risk Engine: Instantly computes fraud probability percentage using trained Random Forest model.",
    "• Model Interpretability: Displays bar chart of top contributing features for the current scoring session.",
    "• How to Run: Execute 'python -m streamlit run app/dashboard.py' (Active at http://localhost:8501)."
])

# --- SLIDE 9: Team Task Allocation ---
slide9 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide9)
add_header(slide9, "Team Task Allocation & Responsibility Matrix")

add_card(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4),
         "1. Lead AI/ML Engineer", [
    "• Model selection & hyperparameter tuning",
    "• Evaluation metrics (ROC-AUC, Precision, Recall)",
    "• Deliverables: train_model.py, fraud_model.joblib"
])

add_card(slide9, Inches(6.8), Inches(1.8), Inches(5.733), Inches(2.4),
         "2. Data & Feature Engineer", [
    "• Synthetic dataset generation & noise modeling",
    "• Feature distribution design (12 signals)",
    "• Deliverables: generate_data.py, synthetic_kyc.csv"
])

add_card(slide9, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.3),
         "3. Full-Stack / UI Engineer", [
    "• Streamlit dashboard application development",
    "• Interactive preset profile loader",
    "• Deliverables: dashboard.py"
])

add_card(slide9, Inches(6.8), Inches(4.5), Inches(5.733), Inches(2.3),
         "4. Research & Documentation", [
    "• Literature survey & project report writing",
    "• Presentation slide deck preparation",
    "• Deliverables: report.md, report.pdf, presentation.pptx"
])

# --- SLIDE 10: Conclusion & Next Steps ---
slide10 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide10)
add_header(slide10, "Conclusion & Future Roadmap")

add_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
         "Key Conclusions", [
    "• Multimodal Detection Works: Combining KYC + behavioral signals achieves high precision (0.95) and recall (0.75).",
    "• Identity Freshness is Primary: Email and phone age are the single strongest predictors of synthetic identity creation.",
    "• Operational Readiness: Streamlit dashboard provides an intuitive interface for live underwriting integration."
])

add_card(slide10, Inches(6.8), Inches(1.8), Inches(5.733), Inches(5.0),
         "Future Roadmap", [
    "• Graph Neural Networks (GNNs): Build identity resolution graphs to link applications across multiple institutions.",
    "• Behavioral Telemetry SDK: Capture real-time keystroke dynamics and mouse movements via JavaScript SDK.",
    "• Production Calibration: Calibrate risk thresholds against real anonymized banking datasets."
])

prs.save(PPT_PATH)
print("Successfully generated PPTX presentation deck at:", PPT_PATH)

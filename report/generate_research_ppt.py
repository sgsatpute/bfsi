import os
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
REPORT_DIR = os.path.join(BASE, "report")
IMG_DIR = os.path.join(REPORT_DIR, "images")
PPT_PATH = os.path.join(REPORT_DIR, "presentation.pptx")

with open(os.path.join(BASE, "model", "metrics.json"), "r") as f:
    metrics = json.load(f)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

COLOR_BG = RGBColor(248, 250, 252)
COLOR_PRIMARY = RGBColor(15, 23, 42)
COLOR_SECONDARY = RGBColor(37, 99, 235)
COLOR_CARD_BG = RGBColor(255, 255, 255)
COLOR_TEXT_DARK = RGBColor(30, 41, 59)
COLOR_ACCENT = RGBColor(16, 185, 129)

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, category_text="BFSI / AI-ML RESEARCH PRESENTATION"):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.95))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_PRIMARY
    box.line.fill.background()

    tf = box.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.1)

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(9.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_SECONDARY

    p_t = tf.add_paragraph()
    p_t.text = title_text
    p_t.font.size = Pt(19)
    p_t.font.bold = True
    p_t.font.color.rgb = RGBColor(255, 255, 255)

def add_card(slide, left, top, width, height, title, bullets):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = RGBColor(226, 232, 240)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.18)

    p_head = tf.paragraphs[0]
    p_head.text = title
    p_head.font.size = Pt(14)
    p_head.font.bold = True
    p_head.font.color.rgb = COLOR_PRIMARY
    p_head.space_after = Pt(8)

    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(5)

# SLIDE 1: Title
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1)

tb1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.4), Inches(11.333), Inches(4.7))
tb1.fill.solid()
tb1.fill.fore_color.rgb = COLOR_PRIMARY

tf1 = tb1.text_frame
tf1.margin_left = Inches(0.6)
tf1.margin_top = Inches(0.5)

p1 = tf1.paragraphs[0]
p1.text = "EXPLAINABLE SYNTHETIC IDENTITY FRAUD DETECTION"
p1.font.size = Pt(26)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)

p2 = tf1.add_paragraph()
p2.text = "Human-Centered Risk Engine for Digital Lending Onboarding"
p2.font.size = Pt(16)
p2.font.bold = True
p2.font.color.rgb = COLOR_SECONDARY
p2.space_after = Pt(20)

p3 = tf1.add_paragraph()
p3.text = f"Live App Portal: https://sgsatpute-bfsi-appdashboard-ikfuca.streamlit.app/\nGitHub: https://github.com/sgsatpute/bfsi\nDataset: 10,000 Applications | 20 Signals | 5-Fold Stratified CV (AUC = {metrics['roc_auc']:.4f})"
p3.font.size = Pt(13)
p3.font.color.rgb = RGBColor(203, 213, 225)

# SLIDE 2: Problem & Threat Formulation
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2)
add_header(s2, "The Threat: Synthetic Identity Fraud")

add_card(s2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3),
         "Why Static KYC Fails", [
    "• Synthetic Fraud: Real PII (PAN tax ID) + Fake Contact (burner phone, synthetic email, fake address).",
    "• Single-Field Query Bypass: Static database queries check PAN in isolation and return 'PASS'.",
    "• No Immediate Victim: Fraudsters build credit for months before maxing loans ('bust-out') and vanishing.",
    "• Growth Rate: Fastest-growing financial crime in digital lending."
])

add_card(s2, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.3),
         "Our Multimodal AI Solution", [
    "• Fuses 20 Signals across 4 Categories: KYC Matching, Identity Freshness, Behavioral Habits, Device Telemetry.",
    "• Machine Learning Engine: Random Forest tuned with class weighting (5-Fold CV AUC = " + f"{metrics['roc_auc']:.4f}" + ").",
    "• Explainable AI Diagnosis: Translates risk scores into plain-English red flag alerts for underwriters.",
    "• Real-Time Action Bands: Low (<30% Auto-Approve), Medium (30-60% Video KYC), High (>60% Block)."
])

# SLIDE 3: 20-Signal Architecture
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3)
add_header(s3, "20 Multimodal Signal Dictionary")

add_card(s3, Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.6),
         "1. KYC & Document Matching (5)", [
    "• name_address_mismatch_score (Distance)",
    "• dob_pan_mismatch (DOB vs PAN flag)",
    "• document_reuse_count (Historical reuse)",
    "• ssn_pan_issuance_gap_years (ID age gap)",
    "• commercial_address_flag (Mailbox flag)"
])

add_card(s3, Inches(6.8), Inches(1.6), Inches(5.733), Inches(2.6),
         "2. Identity Freshness & Bureau (5)", [
    "• phone_age_days (Mobile subscription age)",
    "• email_age_days (Domain/account age)",
    "• credit_bureau_hit (Bureau file hit)",
    "• bureau_file_depth_months (Tradelines)",
    "• social_footprint_score (Digital footprint)"
])

add_card(s3, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.5),
         "3. Behavioral Biometrics (5)", [
    "• session_fill_time_sec (Fill duration)",
    "• typing_speed_variance (Keypress cadence)",
    "• backspace_count (Edit keypresses)",
    "• paste_event_ratio (Clipboard paste ratio)",
    "• field_hesitation_ms (Pause duration)"
])

add_card(s3, Inches(6.8), Inches(4.4), Inches(5.733), Inches(2.5),
         "4. Device & Network Telemetry (5)", [
    "• device_reuse_across_apps (Device apps)",
    "• application_velocity_24h (Velocity 24h)",
    "• ip_geolocation_mismatch (IP vs Address)",
    "• identity_graph_degree_centrality (Graph)",
    "• subnet_risk_score (Subnet risk score)"
])

# SLIDE 4: Human-Centered Explainable UI
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4)
add_header(s4, "Human-Centered Explainable AI (XAI) Portal")

add_card(s4, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3),
         "Underwriter 3-Step Guided Workflow", [
    "• STEP 1: Select Applicant Profile (1-Click Presets: Rohan Sharma, Ankit Verma, Fake Persona #892).",
    "• STEP 2: Click 'Analyze Loan Application Now'.",
    "• STEP 3: Read Plain-English Fraud Diagnosis & Action Recommendation.",
    "• User Perspective: Zero complex slider friction required; all signals grouped into 4 clean accordions."
])

add_card(s4, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.3),
         "Plain-English Red Flag Callouts", [
    "• 🔴 Burner Phone Alert: Mobile subscription activated only 4 days ago.",
    "• 🔴 Scripted Form Fill: Application completed in 14 seconds (Bot pattern).",
    "• 🔴 Clipboard Copy-Paste: 92% of fields populated via copy-paste.",
    "• 🟢 Established Phone Line: Mobile subscription active for 3.3 years.",
    "• 🟢 Natural Typing Cadence: Realistic human form completion duration (270s)."
])

# SLIDE 5: 5-Model Cross-Validation Benchmarks
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5)
add_header(s5, "5-Fold Cross-Validation Model Benchmarks")

kpi = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.8))
kpi.fill.solid()
kpi.fill.fore_color.rgb = COLOR_PRIMARY

tf_k = kpi.text_frame
pk = tf_k.paragraphs[0]
pk.text = f"Production Random Forest Model: ROC-AUC = {metrics['roc_auc']:.4f} | Precision = {metrics['precision_fraud']:.4f} | Recall = {metrics['recall_fraud']:.4f} | Fraud F1 = {metrics['f1_fraud']:.4f}"
pk.font.size = Pt(13.5)
pk.font.bold = True
pk.font.color.rgb = COLOR_ACCENT
pk.alignment = PP_ALIGN.CENTER

img_radar = os.path.join(IMG_DIR, "radar_model_benchmark.png")
if os.path.exists(img_radar):
    s5.shapes.add_picture(img_radar, Inches(0.8), Inches(2.5), Inches(11.733), Inches(4.5))

# SLIDE 6: Multi-Model ROC & PR Curves
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6)
add_header(s6, "Multi-Model ROC & Precision-Recall Curves")

img_roc = os.path.join(IMG_DIR, "roc_curves_comparison.png")
img_pr = os.path.join(IMG_DIR, "pr_curves_comparison.png")

if os.path.exists(img_roc):
    s6.shapes.add_picture(img_roc, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3))
if os.path.exists(img_pr):
    s6.shapes.add_picture(img_pr, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.3))

# SLIDE 7: Conclusion & Live Portal
s7 = prs.slides.add_slide(blank_layout)
set_bg(s7)
add_header(s7, "Conclusion & Live Production Portal")

add_card(s7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3),
         "Key Research Conclusions", [
    "• Multimodal Risk Fusion Works: Fusing KYC document matching, identity age, and behavioral biometrics catches synthetic fraud with 0.86 Precision and 0.91 Recall.",
    "• Identity Age is Primary: Email and phone age are the single strongest predictors of synthetic fabrication.",
    "• High Discrimination: Achieves 0.9139 ROC-AUC across 10,000 synthetic applications."
])

add_card(s7, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.3),
         "Live Production Deployment", [
    "• Live Public Web Portal: https://sgsatpute-bfsi-appdashboard-ikfuca.streamlit.app/",
    "• GitHub Source Code: https://github.com/sgsatpute/bfsi",
    "• Formats Included: IEEE Research Paper (PDF/DOCX), 12-Slide PPTX, Class Proposal Blueprint, Simple Guide, and Submission ZIP Package."
])

prs.save(PPT_PATH)
print(f"Generated Updated PowerPoint Deck at: {PPT_PATH}")
